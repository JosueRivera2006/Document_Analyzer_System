# Ejecutar pip install fastapi uvicorn python-multipart
# Ejecutar pip install sumy nltk
# Ejecutar pip install langdetect
# Ejecutar pip install deep-translator
# Ejecutar pip install wordcloud
# Ejecutar pip install python-docx PyPDF2 python-pptx
import re
import os
import nltk
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse, FileResponse
import tempfile

# Descargas necesarias, se pueden quitar del script, si se ejecuto solamente una vez
nltk.download('punkt')
nltk.download('punkt_tab')
nltk.download('stopwords')

from langdetect import detect
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lex_rank import LexRankSummarizer
from deep_translator import GoogleTranslator
from wordcloud import WordCloud
from nltk.corpus import stopwords

# Librerias para leer documentos tipo docs, pdf, pptx
from docx import Document
import PyPDF2
from pptx import Presentation

app = FastAPI(
    title="API de Resumen de Documentos",
    description="Sube un documento y obtén un resumen automático con análisis de palabras clave.",
    version="1.0.0"
)

# Crear el resumidor
resumidor = LexRankSummarizer()

# Mapeo de idioma a los tokenizadores que acepta Sumy
MAPA_IDIOMAS = {
    'es': 'spanish',
    'en': 'english',
    'fr': 'french',
    'de': 'german',
    'pt': 'portuguese',
}

EXTENSIONES_SOPORTADAS = {"txt", "csv", "md", "pdf", "docx", "pptx"}


# ─────────────────────────────────────────────
# Utilidades internas
# ─────────────────────────────────────────────

def leer_documento(ruta):
    extension = ruta.lower().split('.')[-1] # para saber si el archivo es txt, pdf, pptx, o docs
    # leyendo txt
    if extension == "txt" or extension == "csv" or extension == "md":
        with open(ruta, encoding="utf-8") as f:
            return f.read()
    # leyendo pdf
    elif extension == "pdf":
        texto = ""
        with open(ruta, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for pagina in reader.pages:
                texto += pagina.extract_text() + " "
        return texto
    # leyendo docx
    elif extension == "docx":
        doc = Document(ruta)
        return "\n".join([p.text for p in doc.paragraphs])
    # leyendo pptx
    elif extension == "pptx":
        prs = Presentation(ruta)
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + " "
        return texto
    # extension del documento no soportada
    else:
        raise ValueError("Formato no soportado")


def generar_resumen(texto, nombre_archivo):
    # Eliminar emojis usando regex
    texto = re.sub(r"[^\w\s.,;:¡!¿?áéíóúüñÁÉÍÓÚÜÑ()-]", "", texto)

    if not texto.strip():
        raise ValueError("El documento está vacío o no se pudo extraer texto.")

    # Detectar idioma
    idioma_detectado = detect(texto)
    # Mapeo de idioma a los tokenizadores que acepta Sumy
    tokenizer_idioma = MAPA_IDIOMAS.get(idioma_detectado, 'english')  # inglés por defecto
    stopwords_idioma = set(stopwords.words(MAPA_IDIOMAS.get(idioma_detectado, 'english'))) # palabras que no se toman en cuenta en la nube de palabras respecto a idioma

    # Contar palabras del documento original
    palabras_originales = len(texto.split())

    # Traducir "Resumen de:" al idioma detectado
    traduccion = GoogleTranslator(source='auto', target=idioma_detectado).translate("Resumen de")
    # Crear el título
    titulo = f"{traduccion} {os.path.splitext(nombre_archivo)[0]}"

    # Procesar el documento completo
    parser = PlaintextParser.from_string(texto, Tokenizer(tokenizer_idioma))

    # Generar resumen: seleccionar un número de oraciones
    num_oraciones = max(5, len(parser.document.sentences) // 8)
    resumen = resumidor(parser.document, num_oraciones)

    # Convertir resumen a texto estructurado en párrafos
    oraciones_resumen = [str(oracion) for oracion in resumen]

    parrafos = []
    for i in range(0, len(oraciones_resumen), 3):  # 3 oraciones por párrafo
        parrafo = " ".join(oraciones_resumen[i:i+3])
        parrafos.append(parrafo)

    resumen_final = "\n\n".join(parrafos)

    # Contar palabras del resumen
    palabras_resumen = len(resumen_final.split())

    # Crear objeto WordCloud sin generar imagen
    wordcloud = WordCloud(width=800, height=400, background_color="white", stopwords=stopwords_idioma).generate(resumen_final)

    # Obtener las palabras y sus frecuencias como diccionario
    frecuencias = wordcloud.words_

    # Convertir a lista de palabras más frecuentes (Top 15)
    palabras_clave = list(frecuencias.keys())[:15]

    return {
        "titulo": titulo,
        "idioma_detectado": idioma_detectado,
        "palabras_originales": palabras_originales,
        "palabras_resumen": palabras_resumen,
        "resumen": resumen_final,
        "parrafos": parrafos,
        "palabras_clave": palabras_clave,
    }


# ─────────────────────────────────────────────
# Endpoints
# ─────────────────────────────────────────────

@app.get("/", summary="Estado de la API")
def raiz():
    return {"estado": "activo", "mensaje": "API de Resumen de Documentos lista."}


@app.post("/resumir", summary="Resumir un documento subido")
async def resumir_documento(archivo: UploadFile = File(...)):
    """
    Sube un documento (.txt, .csv, .md, .pdf, .docx, .pptx)
    y recibe un resumen automático con análisis de palabras clave.

    Retorna:
    - titulo: Título del resumen
    - idioma_detectado: Código de idioma (es, en, fr, etc.)
    - palabras_originales: Total de palabras en el documento
    - palabras_resumen: Total de palabras en el resumen
    - resumen: Texto completo del resumen
    - parrafos: Lista de párrafos del resumen
    - palabras_clave: Top 15 palabras más relevantes
    """
    # Nombre del archivo
    nombre = archivo.filename or "documento"
    extension = nombre.lower().split('.')[-1]

    if extension not in EXTENSIONES_SOPORTADAS:
        raise HTTPException(
            status_code=400,
            detail=f"Formato '{extension}' no soportado. Use: {', '.join(EXTENSIONES_SOPORTADAS)}"
        )

    # Guardar temporalmente el archivo para procesarlo
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{extension}") as tmp:
        contenido = await archivo.read()
        tmp.write(contenido)
        ruta_tmp = tmp.name

    try:
        # Leer el texto desde el archivo
        texto = leer_documento(ruta_tmp)
        resultado = generar_resumen(texto, nombre)
    except ValueError as e:
        raise HTTPException(status_code=422, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al procesar el documento: {str(e)}")
    finally:
        os.unlink(ruta_tmp)  # Eliminar archivo temporal

    return JSONResponse(content=resultado)


@app.post("/resumir/texto", summary="Resumir texto plano enviado en el body")
async def resumir_texto(payload: dict):
    """
    Envía texto plano en el body JSON y recibe un resumen.

    Body esperado:
    {
        "texto": "El texto a resumir...",
        "nombre": "mi_documento"
    }
    """
    texto = payload.get("texto", "").strip()
    nombre = payload.get("nombre", "documento")

    if not texto:
        raise HTTPException(status_code=400, detail="El campo 'texto' no puede estar vacío.")

    try:
        resultado = generar_resumen(texto, nombre)
    except ValueError as e:
        raise HTTPException(status_code=422, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al procesar el texto: {str(e)}")

    return JSONResponse(content=resultado)


@app.post("/resumir/exportar", summary="Resumir y descargar como .txt")
async def resumir_y_exportar(archivo: UploadFile = File(...)):
    """
    Igual que /resumir pero devuelve el resumen como archivo descargable (.txt).
    """
    # Nombre del archivo
    nombre = archivo.filename or "documento"
    extension = nombre.lower().split('.')[-1]

    if extension not in EXTENSIONES_SOPORTADAS:
        raise HTTPException(
            status_code=400,
            detail=f"Formato '{extension}' no soportado."
        )

    # Guardar temporalmente el archivo para procesarlo
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{extension}") as tmp:
        contenido = await archivo.read()
        tmp.write(contenido)
        ruta_tmp = tmp.name

    try:
        # Leer el texto desde el archivo
        texto = leer_documento(ruta_tmp)
        resultado = generar_resumen(texto, nombre)
    except ValueError as e:
        raise HTTPException(status_code=422, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        os.unlink(ruta_tmp)  # Eliminar archivo temporal

    # Guardar el resumen en un archivo
    with tempfile.NamedTemporaryFile(
        delete=False, suffix=".txt", mode="w", encoding="utf-8"
    ) as out:
        out.write(resultado["titulo"] + "\n\n")
        out.write(resultado["resumen"])
        ruta_salida = out.name

    nombre_descarga = f"resumen_{os.path.splitext(nombre)[0]}.txt"
    return FileResponse(
        path=ruta_salida,
        filename=nombre_descarga,
        media_type="text/plain"
    )